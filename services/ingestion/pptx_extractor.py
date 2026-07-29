from pptx import Presentation


def extract_pptx(file_path: str) -> list[str]:
    prs = Presentation(file_path)
    chunks = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        slide_text.append(text)
        if slide_text:
            chunks.append(f"Slide {slide_num}: " + " ".join(slide_text))
    return chunks